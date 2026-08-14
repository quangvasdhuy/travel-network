from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ──────────────────────────────────────────────
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # Navy dark bg
C_ACCENT  = RGBColor(0xE9, 0x4F, 0x37)   # Red-orange accent
C_BLUE    = RGBColor(0x16, 0x72, 0xAA)   # Blue accent
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xF0, 0xF4, 0xF8)   # Light bg for content slides
C_GRAY    = RGBColor(0x55, 0x65, 0x6E)   # Body text gray
C_YELLOW  = RGBColor(0xFF, 0xC3, 0x00)   # Highlight yellow
C_GREEN   = RGBColor(0x2E, 0xCC, 0x71)   # Green for code/success
C_CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)  # Code block dark bg
C_CODE_FG = RGBColor(0xA8, 0xFF, 0x78)  # Code green text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helper: add rectangle ──────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, transparency=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

# ── Helper: add textbox ────────────────────────────────────────
def add_text(slide, text, x, y, w, h,
             font_size=20, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox

# ── Helper: add multi-paragraph textbox ───────────────────────
def add_multiline(slide, lines, x, y, w, h,
                  font_size=16, color=C_GRAY, line_spacing=1.15):
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (txt, sz, bold, col, align, indent) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if indent:
            p.level = 1
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = "Segoe UI"
    return txBox

def line(txt, sz=16, bold=False, col=C_GRAY, align=PP_ALIGN.LEFT, indent=False):
    return (txt, sz, bold, col, align, indent)

def heading(txt, sz=22, col=C_DARK):
    return (txt, sz, True, col, PP_ALIGN.LEFT, False)

def sub(txt, sz=15, col=C_GRAY):
    return (txt, sz, False, col, PP_ALIGN.LEFT, True)

def empty():
    return ("", 10, False, C_GRAY, PP_ALIGN.LEFT, False)

# ── Helper: code block ────────────────────────────────────────
def add_code_block(slide, code_text, x, y, w, h):
    bg = add_rect(slide, x, y, w, h, C_CODE_BG)
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                      w - Inches(0.3), h - Inches(0.2))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for codeline in code_text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = codeline
        run.font.size = Pt(11)
        run.font.name = "Consolas"
        run.font.color.rgb = C_CODE_FG

# ── Helper: slide background ──────────────────────────────────
def set_bg(slide, color):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

# ── Helper: section divider bar ───────────────────────────────
def add_divider(slide, y=Inches(1.45)):
    add_rect(slide, Inches(0.6), y, Inches(0.07), Inches(0.45), C_ACCENT)

# ── Helper: slide number ──────────────────────────────────────
def add_slide_number(slide, num):
    add_text(slide, str(num), SLIDE_W - Inches(0.7), SLIDE_H - Inches(0.45),
             Inches(0.5), Inches(0.35), font_size=11, color=C_GRAY, align=PP_ALIGN.RIGHT)

# ── Helper: footer bar ────────────────────────────────────────
def add_footer(slide, text="Nhóm 5 • Cơ sở dữ liệu • Travel Network"):
    add_rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), C_DARK)
    add_text(slide, text,
             Inches(0.5), SLIDE_H - Inches(0.36), Inches(10), Inches(0.32),
             font_size=10, color=RGBColor(0x99, 0xAA, 0xBB), align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════
def make_cover(prs):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_DARK)
    # Left accent strip
    add_rect(slide, 0, 0, Inches(0.5), SLIDE_H, C_ACCENT)
    # Top right decorative circle (big, semi-transparent feel via color)
    add_rect(slide, SLIDE_W - Inches(4.5), Inches(0), Inches(4.5), Inches(4.5),
             RGBColor(0x16, 0x2A, 0x48))
    add_rect(slide, SLIDE_W - Inches(3.2), Inches(0.3), Inches(3.0), Inches(3.0),
             RGBColor(0x1A, 0x35, 0x58))

    # Tag line
    add_text(slide, "MÔN: CƠ SỞ DỮ LIỆU",
             Inches(0.9), Inches(1.1), Inches(9), Inches(0.5),
             font_size=13, color=C_ACCENT, bold=True)
    # Main title
    add_text(slide, "ỨNG DỤNG COUCHBASE",
             Inches(0.9), Inches(1.8), Inches(9), Inches(1.0),
             font_size=44, bold=True, color=C_WHITE)
    add_text(slide, "VÀO DỰ ÁN TRAVEL NETWORK",
             Inches(0.9), Inches(2.7), Inches(9), Inches(0.9),
             font_size=36, bold=True, color=C_YELLOW)
    # Divider line
    add_rect(slide, Inches(0.9), Inches(3.65), Inches(5), Inches(0.05), C_ACCENT)
    # Subtitle
    add_text(slide, "Mạng xã hội du lịch • NoSQL Document Database",
             Inches(0.9), Inches(3.8), Inches(9), Inches(0.5),
             font_size=16, color=RGBColor(0x99, 0xBB, 0xDD))
    # Group info
    add_text(slide, "Nhóm 5",
             Inches(0.9), Inches(5.2), Inches(4), Inches(0.4),
             font_size=15, bold=True, color=C_WHITE)
    add_text(slide, "Học kỳ 2 • 2024–2025",
             Inches(0.9), Inches(5.6), Inches(4), Inches(0.35),
             font_size=13, color=RGBColor(0x88, 0xAA, 0xCC))

make_cover(prs)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — MỤC LỤC
# ══════════════════════════════════════════════════════════════
def make_toc(prs):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_LIGHT)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), C_DARK)
    add_text(slide, "MỤC LỤC",
             Inches(0.6), Inches(0.35), Inches(10), Inches(0.7),
             font_size=30, bold=True, color=C_WHITE)
    add_footer(slide)
    add_slide_number(slide, 2)

    items = [
        ("1", "Giới thiệu dự án"),
        ("2", "Tại sao chọn Couchbase?"),
        ("3", "Kiến trúc cơ sở dữ liệu"),
        ("4", "Thiết kế Document (Data Modeling)"),
        ("5", "Truy vấn dữ liệu với N1QL"),
        ("6", "Index và tối ưu hiệu suất"),
        ("7", "Các kỹ thuật đặc trưng của Couchbase"),
        ("8", "Demo luồng nghiệp vụ thực tế"),
        ("9", "Câu hỏi phản biện thường gặp"),
        ("10", "Kết luận"),
    ]
    cols = [items[:5], items[5:]]
    for ci, col in enumerate(cols):
        cx = Inches(0.7) + ci * Inches(6.5)
        for ri, (num, title) in enumerate(col):
            ry = Inches(1.55) + ri * Inches(0.95)
            add_rect(slide, cx, ry, Inches(0.5), Inches(0.5), C_ACCENT)
            add_text(slide, num, cx, ry + Inches(0.05), Inches(0.5), Inches(0.42),
                     font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, title, cx + Inches(0.6), ry + Inches(0.1),
                     Inches(5.7), Inches(0.42),
                     font_size=16, bold=False, color=C_DARK)

make_toc(prs)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — GIỚI THIỆU DỰ ÁN
# ══════════════════════════════════════════════════════════════
def make_intro(prs):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_LIGHT)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), C_DARK)
    add_divider(slide, Inches(0.9))
    add_text(slide, "1. GIỚI THIỆU DỰ ÁN",
             Inches(0.75), Inches(0.32), Inches(10), Inches(0.7),
             font_size=26, bold=True, color=C_WHITE)
    add_footer(slide); add_slide_number(slide, 3)

    # Left column — description
    add_text(slide, "Travel Network là gì?",
             Inches(0.6), Inches(1.5), Inches(6), Inches(0.45),
             font_size=18, bold=True, color=C_DARK)
    desc = ("Mạng xã hội dành cho người đam mê du lịch, tương tự Instagram "
            "nhưng tập trung vào nội dung du lịch. Người dùng có thể:")
    add_text(slide, desc, Inches(0.6), Inches(2.0), Inches(5.9), Inches(0.8),
             font_size=14, color=C_GRAY)

    features = [
        "Đăng ký tài khoản và tạo hồ sơ cá nhân",
        "Chia sẻ bài đăng kèm hình ảnh về chuyến đi",
        "Tạo và quản lý kế hoạch chuyến đi",
        "Follow người dùng khác và xem feed cá nhân",
        "Khám phá điểm đến trên toàn thế giới",
        "Like, comment và tương tác bài đăng",
    ]
    for i, f in enumerate(features):
        fy = Inches(2.85) + i * Inches(0.56)
        add_rect(slide, Inches(0.6), fy + Inches(0.12), Inches(0.18), Inches(0.18), C_ACCENT)
        add_text(slide, f, Inches(0.88), fy, Inches(5.5), Inches(0.5),
                 font_size=13, color=C_GRAY)

    # Right column — tech stack cards
    add_text(slide, "Stack công nghệ",
             Inches(7.1), Inches(1.5), Inches(5.6), Inches(0.45),
             font_size=18, bold=True, color=C_DARK)
    stack = [
        ("Backend",   "Node.js + Express.js",           C_BLUE),
        ("Database",  "Couchbase (NoSQL)",               C_ACCENT),
        ("Frontend",  "React + Vite + TailwindCSS",      RGBColor(0x06,0xB6,0xD4)),
        ("Storage",   "Cloudinary (ảnh/video)",          RGBColor(0x10,0xB9,0x81)),
        ("Auth",      "JWT — JSON Web Token",            RGBColor(0xF5,0x9E,0x0B)),
    ]
    for i, (label, val, col) in enumerate(stack):
        sy = Inches(2.05) + i * Inches(0.92)
        add_rect(slide, Inches(7.1), sy, Inches(5.6), Inches(0.78), col)
        add_text(slide, label.upper(), Inches(7.2), sy + Inches(0.05),
                 Inches(2), Inches(0.32), font_size=10, bold=True,
                 color=RGBColor(0xFF,0xFF,0xFF))
        add_text(slide, val, Inches(7.2), sy + Inches(0.35),
                 Inches(5.3), Inches(0.38), font_size=15, bold=True, color=C_WHITE)

make_intro(prs)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — TẠI SAO CHỌN COUCHBASE
# ══════════════════════════════════════════════════════════════
def make_why(prs):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_LIGHT)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), C_DARK)
    add_text(slide, "2. TẠI SAO CHỌN COUCHBASE?",
             Inches(0.75), Inches(0.32), Inches(10), Inches(0.7),
             font_size=26, bold=True, color=C_WHITE)
    add_footer(slide); add_slide_number(slide, 4)

    # 4 reason cards
    reasons = [
        (C_ACCENT,                        "Dữ liệu phi cấu trúc",
         "Profile, post, trip... có cấu trúc lồng nhau tự nhiên dạng JSON. Thêm field mới không cần migration."),
        (C_BLUE,                          "Đọc nhiều hơn ghi",
         "Feed, danh sách bài đăng được đọc liên tục. Memory-First của Couchbase cache dữ liệu hot trên RAM."),
        (RGBColor(0x10,0xB9,0x81),        "Scale ngang dễ dàng",
         "Khi người dùng tăng, thêm node vào cluster thay vì nâng cấu hình 1 máy chủ."),
        (RGBColor(0xF5,0x9E,0x0B),        "Key-Value + N1QL",
         "Hỗ trợ cả lookup O(1) theo key (cực nhanh) và N1QL query linh hoạt như SQL."),
    ]
    for i, (col, title, desc) in enumerate(reasons):
        cx = Inches(0.5) + i * Inches(3.2)
        add_rect(slide, cx, Inches(1.5), Inches(3.05), Inches(0.55), col)
        add_text(slide, title, cx + Inches(0.1), Inches(1.53),
                 Inches(2.9), Inches(0.48),
                 font_size=14, bold=True, color=C_WHITE)
        add_rect(slide, cx, Inches(2.05), Inches(3.05), Inches(2.5),
                 RGBColor(0xF8,0xF9,0xFA))
        add_text(slide, desc, cx + Inches(0.12), Inches(2.15),
                 Inches(2.85), Inches(2.3),
                 font_size=12, color=C_DARK)

    # Comparison table
    add_text(slide, "So sánh Couchbase vs SQL",
             Inches(0.5), Inches(4.75), Inches(6), Inches(0.4),
             font_size=16, bold=True, color=C_DARK)

    headers = ["Tiêu chí", "SQL (MySQL)", "Couchbase"]
    rows = [
        ["Cấu trúc dữ liệu", "Schema cố định, migration khi thay đổi", "Document JSON linh hoạt"],
        ["Dữ liệu lồng nhau", "Tách nhiều bảng + JOIN", "Nhúng thẳng vào document"],
        ["Scale", "Vertical (nâng RAM/CPU)", "Horizontal (thêm node)"],
        ["Tốc độ đọc theo key", "Index lookup", "O(1) Key-Value"],
        ["Truy vấn", "SQL", "N1QL (gần giống SQL)"],
    ]
    col_w = [Inches(2.8), Inches(4.2), Inches(3.0)]
    col_x = [Inches(0.5), Inches(3.35), Inches(7.6)]
    hy = Inches(5.2)
    for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
        bg = C_DARK if ci == 0 else (C_BLUE if ci == 2 else RGBColor(0x44,0x55,0x66))
        add_rect(slide, cx, hy, cw - Inches(0.05), Inches(0.38), bg)
        add_text(slide, hdr, cx + Inches(0.08), hy + Inches(0.04),
                 cw - Inches(0.15), Inches(0.32),
                 font_size=12, bold=True, color=C_WHITE)
    for ri, row in enumerate(rows):
        ry = Inches(5.58) + ri * Inches(0.34)
        row_bg = RGBColor(0xF0,0xF4,0xF8) if ri % 2 == 0 else C_WHITE
        for ci, (cell, cw, cx) in enumerate(zip(row, col_w, col_x)):
            add_rect(slide, cx, ry, cw - Inches(0.05), Inches(0.32), row_bg)
            fc = C_ACCENT if ci == 2 else C_DARK
            add_text(slide, cell, cx + Inches(0.07), ry + Inches(0.03),
                     cw - Inches(0.15), Inches(0.28), font_size=11, color=fc)

make_why(prs)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — KIẾN TRÚC DATABASE
# ══════════════════════════════════════════════════════════════
def make_arch(prs):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_LIGHT)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), C_DARK)
    add_text(slide, "3. KIẾN TRÚC CƠ SỞ DỮ LIỆU",
             Inches(0.75), Inches(0.32), Inches(10), Inches(0.7),
             font_size=26, bold=True, color=C_WHITE)
    add_footer(slide); add_slide_number(slide, 5)

    # 4 bucket cards
    buckets = [
        ("travel_users",   "Tài khoản & hồ sơ người dùng",
         ["• user document", "• email, username, passwordHash",
          "• profile (ảnh, bio, location)", "• interests, stats, preferences"],
         C_BLUE),
        ("travel_content", "Bài đăng (Posts)",
         ["• post document", "• content (text, media)", "• likes, comments (nhúng)",
          "• tags, location, visibility"],
         C_ACCENT),
        ("travel_trips",   "Chuyến đi & Điểm đến",
         ["• trip document", "• destination document",
          "• startDate, endDate, status", "• danh sách destinations"],
         RGBColor(0x10,0xB9,0x81)),
        ("travel_social",  "Quan hệ Follow",
         ["• connection document", "• followerId → followingId",
          "• composite key lookup O(1)", "• status: active/blocked"],
         RGBColor(0xF5,0x9E,0x0B)),
    ]
    for i, (name, title, details, col) in enumerate(buckets):
        bx = Inches(0.35) + i * Inches(3.22)
        add_rect(slide, bx, Inches(1.45), Inches(3.1), Inches(0.55), col)
        add_text(slide, name, bx + Inches(0.1), Inches(1.48),
                 Inches(2.95), Inches(0.48), font_size=13, bold=True, color=C_WHITE)
        add_rect(slide, bx, Inches(2.0), Inches(3.1), Inches(3.6),
                 RGBColor(0xFF,0xFF,0xFF))
        add_text(slide, title, bx + Inches(0.1), Inches(2.05),
                 Inches(2.95), Inches(0.38), font_size=12, bold=True, color=col)
        for j, d in enumerate(details):
            add_text(slide, d, bx + Inches(0.1), Inches(2.48) + j * Inches(0.58),
                     Inches(2.95), Inches(0.52), font_size=11, color=C_GRAY)

    # SQL mapping table at bottom
    add_text(slide, "Mapping với SQL",
             Inches(0.35), Inches(5.75), Inches(4), Inches(0.35),
             font_size=14, bold=True, color=C_DARK)
    mapping = [
        ("Cluster", "Database Server"),
        ("Bucket", "Database / Schema"),
        ("Collection", "Table (schema-less)"),
        ("Document", "Row (JSON)"),
        ("Document Key", "Primary Key"),
    ]
    for i, (cb, sql) in enumerate(mapping):
        mx = Inches(0.35) + i * Inches(2.58)
        add_rect(slide, mx, Inches(6.15), Inches(2.5), Inches(0.32), C_DARK)
        add_text(slide, cb, mx + Inches(0.07), Inches(6.17), Inches(1.2), Inches(0.28),
                 font_size=10, bold=True, color=C_YELLOW)
        add_text(slide, "→ " + sql, mx + Inches(1.25), Inches(6.17), Inches(1.2), Inches(0.28),
                 font_size=10, color=C_WHITE)

make_arch(prs)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — DATA MODELING: USER
# ══════════════════════════════════════════════════════════════
def make_model_user(prs):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_LIGHT)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), C_DARK)
    add_text(slide, "4. THIẾT KẾ DOCUMENT — User",
             Inches(0.75), Inches(0.32), Inches(10), Inches(0.7),
             font_size=26, bold=True, color=C_WHITE)
    add_footer(slide); add_slide_number(slide, 6)

    add_text(slide, 'Document Key: user::{uuid}   |   Bucket: travel_users',
             Inches(0.6), Inches(1.38), Inches(8), Inches(0.35),
             font_size=12, color=C_BLUE, bold=True)

    user_json = '''{
  "id": "550e8400-e29b-41d4-a716-446655440000",
  "type": "user",
  "email": "sarah@example.com",
  "username": "traveler_sarah",
  "passwordHash": "$2a$10$...",
  "profile": {
    "firstName": "Sarah",  "lastName": "Johnson",
    "bio": "Adventure seeker | 30+ countries",
    "profilePhoto": "https://cloudinary.com/...",
    "location": { "city": "San Francisco", "country": "USA" }
  },
  "interests": ["hiking", "photography", "backpacking"],
  "stats": { "tripCount": 5, "postCount": 23,
             "followerCount": 142, "followingCount": 87 },
  "status": "active",
  "createdAt": "2024-01-01T00:00:00.000Z"
}'''
    add_code_block(slide, user_json, Inches(0.5), Inches(1.8), Inches(6.5), Inches(5.3))

    # Notes on right
    notes = [
        ("Điểm đáng chú ý:", 15, True, C_DARK),
        ("", 8, False, C_GRAY),
        ('Field "type: user"', 13, True, C_ACCENT),
        ("Phân biệt loại document trong cùng 1 bucket. Kết hợp Partial Index để query nhanh.", 12, False, C_GRAY),
        ("", 8, False, C_GRAY),
        ("Mảng interests", 13, True, C_ACCENT),
        ("Nhúng thẳng vào document, không cần bảng quan hệ trung gian như SQL.", 12, False, C_GRAY),
        ("", 8, False, C_GRAY),
        ("stats (denormalize)", 13, True, C_ACCENT),
        ("Lưu sẵn số đếm (followerCount...) để tránh COUNT query tốn kém mỗi lần hiển thị.", 12, False, C_GRAY),
        ("", 8, False, C_GRAY),
        ("Nested objects", 13, True, C_ACCENT),
        ("profile.location là object lồng trong document, N1QL có thể index sâu vào field này.", 12, False, C_GRAY),
    ]
    add_multiline(slide, [line(t, s, b, c) for t, s, b, c in notes],
                  Inches(7.2), Inches(1.8), Inches(5.8), Inches(5.3))

make_model_user(prs)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — DATA MODELING: POST
# ══════════════════════════════════════════════════════════════
def make_model_post(prs):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_LIGHT)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), C_DARK)
    add_text(slide, "4. THIẾT KẾ DOCUMENT — Post & Connection",
             Inches(0.75), Inches(0.32), Inches(11), Inches(0.7),
             font_size=26, bold=True, color=C_WHITE)
    add_footer(slide); add_slide_number(slide, 7)

    # Post JSON
    add_text(slide, 'post::{uuid}  →  travel_content',
             Inches(0.5), Inches(1.38), Inches(6), Inches(0.3),
             font_size=11, color=C_BLUE, bold=True)
    post_json = '''{
  "id": "post-uuid",  "type": "post",
  "authorId": "user-uuid",
  "authorUsername": "traveler_sarah",
  "authorPhoto": "https://...",
  "content": {
    "text": "Sunrise at Machu Picchu!",
    "media": [{"type":"image","url":"https://..."}]
  },
  "location": {"name":"Machu Picchu","coordinates":
               {"lat":-13.16,"lon":-72.54}},
  "tags": ["travel","peru","adventure"],
  "stats": {"likeCount":42,"commentCount":8},
  "interactions": {
    "likes": ["user-uuid-1","user-uuid-2"],
    "comments": [{"id":"c-uuid","userId":"uid1",
      "username":"explorer_mike",
      "text":"Amazing!"}]
  },
  "visibility": "public"
}'''
    add_code_block(slide, post_json, Inches(0.5), Inches(1.75), Inches(6.0), Inches(5.35))

    # Connection JSON
    add_text(slide, 'connection::{follower}::{following}  →  travel_social',
             Inches(6.85), Inches(1.38), Inches(6.1), Inches(0.3),
             font_size=11, color=C_BLUE, bold=True)
    conn_json = '''{
  "type": "connection",
  "followerId": "user-uuid-A",
  "followerUsername": "traveler_sarah",
  "followingId": "user-uuid-B",
  "followingUsername": "explorer_mike",
  "status": "active",
  "createdAt": "2024-01-10T00:00:00Z"
}'''
    add_code_block(slide, conn_json, Inches(6.85), Inches(1.75), Inches(6.1), Inches(2.7))

    # Key insight boxes
    insights = [
        (C_ACCENT,              "Embedded Document",
         "Comments nhúng vào post — luôn đọc cùng nhau, không cần JOIN riêng."),
        (C_BLUE,                "Denormalization",
         "authorUsername & authorPhoto lưu trong post để hiển thị feed không cần JOIN."),
        (RGBColor(0x10,0xB9,0x81), "Composite Key O(1)",
         "connection::A::B → kiểm tra A có follow B chỉ bằng 1 Key-Value lookup."),
    ]
    for i, (col, title, desc) in enumerate(insights):
        iy = Inches(4.6) + i * Inches(0.9)
        add_rect(slide, Inches(6.85), iy, Inches(0.12), Inches(0.65), col)
        add_text(slide, title, Inches(7.1), iy, Inches(5.7), Inches(0.32),
                 font_size=13, bold=True, color=col)
        add_text(slide, desc, Inches(7.1), iy + Inches(0.3), Inches(5.7), Inches(0.5),
                 font_size=11, color=C_GRAY)

make_model_post(prs)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — N1QL QUERIES (part 1)
# ══════════════════════════════════════════════════════════════
def make_n1ql_1(prs):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_LIGHT)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), C_DARK)
    add_text(slide, "5. TRUY VẤN DỮ LIỆU VỚI N1QL  (1/2)",
             Inches(0.75), Inches(0.32), Inches(11), Inches(0.7),
             font_size=26, bold=True, color=C_WHITE)
    add_footer(slide); add_slide_number(slide, 8)

    add_text(slide,
             "N1QL — Non-First Normal Form Query Language. Cú pháp gần giống SQL nhưng query trên JSON document.",
             Inches(0.6), Inches(1.38), Inches(12), Inches(0.35),
             font_size=12, color=C_GRAY, italic=True)

    # Q1
    add_rect(slide, Inches(0.5), Inches(1.8), Inches(0.35), Inches(0.32), C_ACCENT)
    add_text(slide, "Q1", Inches(0.5), Inches(1.82), Inches(0.35), Inches(0.28),
             font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Tìm user theo email (Named Parameter — tránh SQL Injection)",
             Inches(0.93), Inches(1.82), Inches(7), Inches(0.28),
             font_size=12, color=C_DARK)
    q1 = ("SELECT META().id, u.*\n"
          "FROM travel_users u\n"
          "WHERE u.type = 'user' AND u.email = $email\n"
          "LIMIT 1")
    add_code_block(slide, q1, Inches(0.5), Inches(2.18), Inches(6.1), Inches(1.32))

    # Q2
    add_rect(slide, Inches(6.85), Inches(1.8), Inches(0.35), Inches(0.32), C_BLUE)
    add_text(slide, "Q2", Inches(6.85), Inches(1.82), Inches(0.35), Inches(0.28),
             font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Bài đăng phổ biến nhất (tính toán inline)",
             Inches(7.28), Inches(1.82), Inches(5.5), Inches(0.28),
             font_size=12, color=C_DARK)
    q2 = ("SELECT META().id, p.*\n"
          "FROM travel_content p\n"
          "WHERE p.type = 'post' AND p.visibility = 'public'\n"
          "ORDER BY (\n"
          "  IFMISSINGORNULL(p.stats.likeCount, 0) +\n"
          "  IFMISSINGORNULL(p.stats.commentCount, 0)\n"
          ") DESC, p.createdAt DESC\n"
          "LIMIT $limit OFFSET $offset")
    add_code_block(slide, q2, Inches(6.85), Inches(2.18), Inches(6.1), Inches(1.9))

    # Q3 — Feed (subquery)
    add_rect(slide, Inches(0.5), Inches(3.6), Inches(0.35), Inches(0.32), RGBColor(0x10,0xB9,0x81))
    add_text(slide, "Q3", Inches(0.5), Inches(3.62), Inches(0.35), Inches(0.28),
             font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Feed cá nhân — Subquery lấy posts từ người đang follow",
             Inches(0.93), Inches(3.62), Inches(11.5), Inches(0.28),
             font_size=12, color=C_DARK)
    q3 = ("SELECT META(p).id, p.*\n"
          "FROM travel_content p\n"
          "WHERE p.type = 'post' AND p.visibility = 'public'\n"
          "  AND p.authorId IN (\n"
          "    SELECT RAW c.followingId\n"
          "    FROM travel_social c\n"
          "    WHERE c.type = 'connection'\n"
          "      AND c.followerId = $userId AND c.status = 'active'\n"
          "  )\n"
          "ORDER BY p.createdAt DESC\n"
          "LIMIT $limit OFFSET $offset")
    add_code_block(slide, q3, Inches(0.5), Inches(3.98), Inches(12.45), Inches(2.65))

make_n1ql_1(prs)
